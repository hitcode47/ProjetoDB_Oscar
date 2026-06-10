"""Gera apresentação .pptx — Projeto Final BD — Oscar AMPAS"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Paleta ────────────────────────────────────────────────────────────
OURO      = RGBColor(0xC8, 0xA0, 0x00)   # dourado Oscar
PRETO     = RGBColor(0x1A, 0x1A, 0x1A)
BRANCO    = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_FG  = RGBColor(0xF5, 0xF5, 0xF5)
CINZA_SUB = RGBColor(0x55, 0x55, 0x55)
AZUL      = RGBColor(0x1F, 0x4E, 0x79)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]   # layout em branco


# ── Helpers ────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=20, bold=False, color=BRANCO,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_bullet_box(slide, items, l, t, w, h,
                   font_size=16, color=PRETO, spacing_pt=6):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(spacing_pt)
        p.level = 0
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txb


def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(blank)
    # fundo preto
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PRETO)
    # faixa dourada vertical esquerda
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, OURO)
    # faixa dourada horizontal inferior
    add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), OURO)

    add_text(slide, title,
             Inches(0.7), Inches(2.2), Inches(12), Inches(2),
             font_size=40, bold=True, color=OURO, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle,
             Inches(0.7), Inches(4.3), Inches(12), Inches(1.2),
             font_size=18, color=BRANCO, align=PP_ALIGN.CENTER)
    return slide


def section_slide(prs, number, title):
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, AZUL)
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, OURO)
    add_text(slide, number,
             Inches(0.7), Inches(2.5), Inches(12), Inches(1),
             font_size=52, bold=True, color=OURO, align=PP_ALIGN.CENTER)
    add_text(slide, title,
             Inches(0.7), Inches(3.6), Inches(12), Inches(1),
             font_size=26, bold=False, color=BRANCO, align=PP_ALIGN.CENTER)
    return slide


def content_slide(prs, title, bullets, highlight=None):
    """Slide com título + lista de bullets. highlight = texto no rodapé."""
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CINZA_FG)
    # barra superior
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), PRETO)
    # linha dourada sob barra
    add_rect(slide, 0, Inches(1.0), SLIDE_W, Inches(0.06), OURO)

    add_text(slide, title,
             Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.75),
             font_size=22, bold=True, color=OURO, align=PP_ALIGN.LEFT)

    add_bullet_box(slide, bullets,
                   Inches(0.6), Inches(1.25), Inches(12.1), Inches(5.5),
                   font_size=17, color=PRETO)

    if highlight:
        add_rect(slide, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7), OURO)
        add_text(slide, highlight,
                 Inches(0.75), Inches(6.52), Inches(11.8), Inches(0.65),
                 font_size=14, bold=True, color=PRETO, align=PP_ALIGN.CENTER)
    return slide


def two_col_slide(prs, title, left_items, right_items,
                  left_title="", right_title=""):
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CINZA_FG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), PRETO)
    add_rect(slide, 0, Inches(1.0), SLIDE_W, Inches(0.06), OURO)

    add_text(slide, title,
             Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.75),
             font_size=22, bold=True, color=OURO)

    mid = Inches(6.75)

    if left_title:
        add_text(slide, left_title,
                 Inches(0.5), Inches(1.2), Inches(5.9), Inches(0.5),
                 font_size=14, bold=True, color=AZUL)
    add_bullet_box(slide, left_items,
                   Inches(0.5), Inches(1.2) + (Pt(18) if left_title else 0),
                   Inches(5.9), Inches(5.5),
                   font_size=15, color=PRETO)

    # divisória
    add_rect(slide, mid, Inches(1.1), Inches(0.04), Inches(6.2), OURO)

    if right_title:
        add_text(slide, right_title,
                 Inches(7.0), Inches(1.2), Inches(5.9), Inches(0.5),
                 font_size=14, bold=True, color=AZUL)
    add_bullet_box(slide, right_items,
                   Inches(7.0), Inches(1.2) + (Pt(18) if right_title else 0),
                   Inches(5.9), Inches(5.5),
                   font_size=15, color=PRETO)
    return slide


# ══════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════

# 1 — Capa
title_slide(prs,
    "Oscar AMPAS — Winner Demographics",
    "Projeto Final de Banco de Dados  |  DCC/UFMG  |  Prof. Pedro H. Barros\n2026")

# 2 — Agenda
content_slide(prs, "Agenda",
    ["1  Dataset e Motivação",
     "2  Modelagem Entidade-Relacionamento",
     "3  Modelo Relacional e Normalização",
     "4  Implementação Física (DDL + PostgreSQL)",
     "5  Pipeline de Carga e Validação",
     "6  Análise Exploratória — 5 Perguntas",
     "7  Principais Achados",
     "8  Conclusão"])

# 3 — Dataset
content_slide(prs, "1  Dataset — Oscar AMPAS Winner Demographics",
    ["Fonte: FiveThirtyEight / AMPAS (Academy of Motion Picture Arts and Sciences)",
     "415 registros  ·  10 colunas  ·  período 1927–2014",
     "5 categorias: Best Actor, Best Actress, Best Supporting Actor,",
     "               Best Supporting Actress, Best Director",
     "87 edições  ·  348 vencedores únicos  ·  335 filmes únicos",
     "",
     "Desafios de qualidade:",
     "  → religion: 62% ausente (estrutural — dado não coletado historicamente)",
     "  → sexual_orientation: string Na tratada como NULL na carga",
     "  → birth_date: 1 registro ausente; birth_year preservado"])

# 4 — Modelo ER
content_slide(prs, "2  Modelo Entidade-Relacionamento",
    ["5 entidades: VENCEDOR · FILME · CATEGORIA · EDICAO · PREMIO",
     "",
     "PREMIO é entidade associativa — conecta as 4 dimensões",
     "  → id_vencedor  FK → VENCEDOR",
     "  → id_filme     FK → FILME",
     "  → id_categoria FK → CATEGORIA",
     "  → id_edicao    FK → EDICAO",
     "",
     "Cardinalidades: todos os relacionamentos N:1 de PREMIO para as dimensões",
     "Participação total em PREMIO (todas as FKs NOT NULL)",
     "",
     "Diagrama gerado via dbdiagram.io  (arquivo: diagrama_er.dbml)"],
    highlight="Um vencedor pode ter múltiplos prêmios — ex: Katharine Hepburn (4 Oscars)")

# 5 — Modelo Relacional
content_slide(prs, "3  Modelo Relacional e Normalização",
    ["vencedor (id_vencedor PK, nome, ano_nascimento, data_nascimento,",
     "          local_nascimento, etnia, religiao, orient_sexual)",
     "filme    (id_filme PK, titulo)",
     "categoria (id_categoria PK, nome)",
     "edicao   (id_edicao PK, ano)",
     "premio   (id_premio PK, id_vencedor FK, id_filme FK,",
     "          id_categoria FK, id_edicao FK)",
     "",
     "Normalização:",
     "  1FN ✔  Valores atômicos em todas as colunas",
     "  2FN ✔  PKs simples — dependências parciais impossíveis",
     "  3FN ✔  Sem dependências transitivas; ano_nascimento mantido pois",
     "          1 registro tem birth_year mas não birth_date completa"])

# 6 — Implementação Física
content_slide(prs, "4  Implementação Física",
    ["SGBD: PostgreSQL 16  |  Container Docker: oscar-db  |  BD: oscar",
     "",
     "Constraints implementadas:",
     "  → PKs (SERIAL) em todas as tabelas",
     "  → NOT NULL em campos obrigatórios",
     "  → ano >= 1927 na tabela edicao",
     "  → FKs em premio (ON DELETE RESTRICT)",
     "  → UNIQUE (id_vencedor, id_categoria, id_edicao) em premio",
     "",
     "Índices criados:",
     "  idx_premio_vencedor  ·  idx_premio_filme",
     "  idx_premio_categoria ·  idx_premio_edicao"])

# 7 — Carga
content_slide(prs, "5  Pipeline de Carga e Validação",
    ["Pipeline Python + psycopg2  (arquivo: 02_carga_dados.py)",
     "",
     "Ordem de inserção (respeita FKs):",
     "  1. edicao     →  87 linhas",
     "  2. categoria  →   5 linhas",
     "  3. filme      → 335 linhas",
     "  4. vencedor   → 348 linhas  (deduplicados por nome)",
     "  5. premio     → 415 linhas",
     "",
     "Validação pós-carga (03_validacao_carga.sql):",
     "  ✔ Zero registros órfãos em todas as FKs",
     "  ✔ Nenhuma duplicata na constraint uq_premio",
     "  ✔ NULLs distribuídos como esperado (religion: 219 · orient_sexual: 10)",
     "  ✔ 30 edições com ≠ 5 prêmios — estrutural (décadas iniciais)"])

# 8 — Perguntas
content_slide(prs, "6  Análise Exploratória — 5 Perguntas Investigativas",
    ["P1  Filmes com múltiplas vitórias na mesma edição",
     "     Técnica: JOIN + GROUP BY + HAVING + STRING_AGG",
     "",
     "P2  Evolução da diversidade étnica por década",
     "     Técnica: GROUP BY com cálculo de década + percentual",
     "",
     "P3  Idade dos vencedores por categoria",
     "     Técnica: AVG / MIN / MAX + UNION ALL",
     "",
     "P4  Maior intervalo entre vitórias consecutivas",
     "     Técnica: CTE + window function LAG",
     "",
     "P5  Primeiro vencedor não-branco por categoria",
     "     Técnica: CTE + ROW_NUMBER + filtro"])

# 9 — Resultados P1 e P2
two_col_slide(prs,
    "7  Principais Achados",
    left_title="P1 — Filmes mais dominantes",
    left_items=[
        "It Happened One Night (1935)",
        "  → 3 categorias na mesma cerimônia",
        "Gone with the Wind (1940)",
        "  → 3 categorias",
        "Going My Way (1945)",
        "  → 3 categorias",
        "",
        "60 filmes no total venceram",
        "2 ou mais categorias na mesma edição"],
    right_title="P2 — Diversidade étnica",
    right_items=[
        "Até os anos 1970: praticamente",
        "  0% de vencedores não-brancos",
        "",
        "Anos 2000: subiu para ~19%",
        "",
        "Ainda longe de representar",
        "a diversidade americana"])

# 10 — Resultados P3 e P4
two_col_slide(prs,
    "7  Principais Achados (cont.)",
    left_title="P3 — Idades extremas",
    left_items=[
        "Mais jovem:",
        "  Tatum O'Neal — 11 anos (1974)",
        "  Best Supporting Actress",
        "",
        "Mais velho:",
        "  Christopher Plummer — 83 anos (2012)",
        "  Best Supporting Actor",
        "",
        "Média mais jovem: Best Actress (37)",
        "Média mais velha: Best Supp. Actor (51)"],
    right_title="P4 — Maior intervalo",
    right_items=[
        "Helen Hayes",
        "  1932 → 1971  (39 anos)",
        "",
        "Katharine Hepburn",
        "  1934 → 1968  (34 anos)",
        "",
        "Revelam carreiras longas e",
        "retornos históricos ao palco",
        "do Oscar"])

# 11 — Resultados P5
content_slide(prs, "7  Principais Achados — P5: Primeira vitória não-branca por categoria",
    ["Best Supporting Actress  →  Hattie McDaniel  (Black)          1940",
     "Best Actor               →  Jose Ferrer       (Hispanic)       1951",
     "Best Supporting Actor    →  Anthony Quinn     (Hispanic)       1953",
     "Best Actress             →  Halle Berry        (Multiracial)   2002",
     "Best Director            →  Ang Lee            (Asian)         2006",
     "",
     "→ Best Director demorou 79 anos para ter um vencedor não-branco",
     "→ Best Actress ficou 75 anos sem diversidade racial",
     "→ O dataset cobre até 2014 — padrões recentes (#OscarsSoWhite, 2016)"],
    highlight="Técnica: CTE + ROW_NUMBER OVER (PARTITION BY categoria ORDER BY ano)  —  window function")

# 12 — Conclusão
content_slide(prs, "8  Conclusão",
    ["Ciclo completo implementado:",
     "  ER  →  Relacional  →  Normalização 3FN  →  DDL  →  Carga  →  EDA",
     "",
     "Decisões técnicas relevantes:",
     "  → Entidade associativa PREMIO captura a granularidade do dataset",
     "  → UNIQUE (vencedor, categoria, edição) previne duplicatas semânticas",
     "  → ano_nascimento mantido pois 1 registro não possui data completa",
     "  → 'Na' tratado como NULL na carga (não string vazia)",
     "",
     "Achado mais impactante:",
     "  → A Academia levou quase 80 anos para premiar um",
     "    diretor não-branco (Ang Lee, 2006)"])

# 13 — Arquivos
content_slide(prs, "Arquivos do Projeto",
    ["Documentacao/",
     "  Etapa0_Exploracao/  01_exploracao.ipynb",
     "  Etapa1_ModeloER/    diagrama_er.dbml",
     "  Etapa3_SQL/         01_ddl_criar_tabelas.sql",
     "                      02_carga_dados.py",
     "                      03_validacao_carga.sql",
     "  Etapa4_EDA/         04_perguntas_investigativas.md",
     "                      05_consultas_sql.sql",
     "                      06_visualizacoes.ipynb",
     "",
     "Aprendizado/",
     "  00_rota_do_projeto.md",
     "  01_entidades_regras.md",
     "  02_modelo_relacional_normalizacao.md",
     "  relatorio_tecnico.md"])

# 14 — Encerramento
title_slide(prs,
    "Obrigado!",
    "Dúvidas?\n\nOscar AMPAS  ·  DCC/UFMG  ·  2026")


# ── Salvar ─────────────────────────────────────────────────────────────
OUT = r"c:\Users\bruno\OneDrive\Área de Trabalho\Projetos\projetoFinalDB\Documentacao\apresentacao_oscar.pptx"
prs.save(OUT)
print(f"Salvo em: {OUT}")
